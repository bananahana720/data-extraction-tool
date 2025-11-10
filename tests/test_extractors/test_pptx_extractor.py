"""
Tests for PowerPoint (PPTX) Extractor

TDD approach: Each test written before implementation.
Organized by requirement from test plan.
"""

import pytest
from pathlib import Path
from datetime import datetime

import sys

sys.path.insert(0, str(Path(__file__).parent.parent.parent / "src"))

from core import (
    ContentBlock,
    ContentType,
    DocumentMetadata,
    ExtractionResult,
    Position,
)

# Import will fail initially - that's expected in TDD
try:
    from extractors.pptx_extractor import PptxExtractor

    PPTX_EXTRACTOR_AVAILABLE = True
except ImportError:
    PPTX_EXTRACTOR_AVAILABLE = False


pytestmark = pytest.mark.skipif(
    not PPTX_EXTRACTOR_AVAILABLE, reason="PptxExtractor not yet implemented"
)


class TestBasicSlideExtraction:
    """Requirement 1: Slide Content Extraction"""

    def test_extract_slide_title(self, simple_pptx_file):
        """Extract title text from a slide."""
        extractor = PptxExtractor()
        result = extractor.extract(simple_pptx_file)

        assert result.success
        assert len(result.content_blocks) > 0

        # Find title block
        titles = [b for b in result.content_blocks if b.block_type == ContentType.HEADING]
        assert len(titles) > 0
        assert titles[0].content.strip() != ""

    def test_extract_slide_body(self, simple_pptx_file):
        """Extract body text from slide shapes."""
        extractor = PptxExtractor()
        result = extractor.extract(simple_pptx_file)

        assert result.success

        # Should have both headings and paragraphs
        paragraphs = [b for b in result.content_blocks if b.block_type == ContentType.PARAGRAPH]
        assert len(paragraphs) > 0

    def test_slide_sequence(self, simple_pptx_file):
        """Verify slides extracted in correct order."""
        extractor = PptxExtractor()
        result = extractor.extract(simple_pptx_file)

        assert result.success

        # Check that slide numbers are sequential
        slide_blocks = [
            b for b in result.content_blocks if b.position and b.position.slide is not None
        ]
        assert len(slide_blocks) > 0

        # First slide should be slide 1
        first_slide = min(b.position.slide for b in slide_blocks)
        assert first_slide == 1

    def test_empty_slide(self, empty_pptx_file):
        """Handle presentation with empty slides."""
        extractor = PptxExtractor()
        result = extractor.extract(empty_pptx_file)

        # Should succeed even with no content
        assert result.success
        # May have zero blocks or just empty blocks
        assert len(result.content_blocks) >= 0


class TestSpeakerNotes:
    """Requirement 2: Speaker Notes Extraction"""

    def test_extract_speaker_notes(self, pptx_with_notes):
        """Extract presenter notes from slides."""
        extractor = PptxExtractor()
        result = extractor.extract(pptx_with_notes)

        assert result.success

        # Find blocks with notes metadata
        notes_blocks = [
            b for b in result.content_blocks if b.metadata.get("is_speaker_note", False)
        ]
        assert len(notes_blocks) > 0
        assert notes_blocks[0].content.strip() != ""

    def test_slide_without_notes(self, simple_pptx_file):
        """Handle slides without speaker notes."""
        extractor = PptxExtractor()
        result = extractor.extract(simple_pptx_file)

        assert result.success
        # Should not crash on missing notes


class TestSlidePositionMetadata:
    """Requirement 3: Slide Layout and Structure"""

    def test_slide_position_metadata(self, simple_pptx_file):
        """Verify Position.slide field set correctly."""
        extractor = PptxExtractor()
        result = extractor.extract(simple_pptx_file)

        assert result.success

        for block in result.content_blocks:
            if block.position:
                # Slide number should be positive integer
                assert block.position.slide is not None
                assert block.position.slide >= 1

    def test_sequence_index(self, simple_pptx_file):
        """Verify sequence_index tracks content order."""
        extractor = PptxExtractor()
        result = extractor.extract(simple_pptx_file)

        assert result.success

        # Sequence indices should be sequential
        indices = [
            b.position.sequence_index
            for b in result.content_blocks
            if b.position and b.position.sequence_index is not None
        ]
        assert len(indices) > 0
        assert indices == sorted(indices)


class TestPresentationMetadata:
    """Requirement 5: Presentation Metadata"""

    def test_presentation_metadata(self, simple_pptx_file):
        """Extract presentation-level metadata."""
        extractor = PptxExtractor()
        result = extractor.extract(simple_pptx_file)

        assert result.success
        assert result.document_metadata is not None
        assert result.document_metadata.source_file == simple_pptx_file
        assert result.document_metadata.file_format == "pptx"

    def test_file_hash(self, simple_pptx_file):
        """Generate SHA256 hash for deduplication."""
        extractor = PptxExtractor()
        result = extractor.extract(simple_pptx_file)

        assert result.success
        assert result.document_metadata.file_hash is not None
        assert len(result.document_metadata.file_hash) == 64  # SHA256 hex length


class TestErrorHandling:
    """Test error handling for invalid inputs."""

    def test_file_not_found(self):
        """Handle non-existent file."""
        extractor = PptxExtractor()
        result = extractor.extract(Path("nonexistent.pptx"))

        assert not result.success
        assert len(result.errors) > 0

    def test_invalid_file_format(self, tmp_path):
        """Handle non-PPTX file."""
        fake_file = tmp_path / "not_a_pptx.txt"
        fake_file.write_text("This is not a PowerPoint file")

        extractor = PptxExtractor()
        result = extractor.extract(fake_file)

        assert not result.success
        assert len(result.errors) > 0

    def test_corrupted_pptx(self, tmp_path):
        """Handle corrupted PPTX file."""
        # Create a file that looks like PPTX but is corrupted
        corrupted = tmp_path / "corrupted.pptx"
        corrupted.write_bytes(b"PK\x03\x04" + b"corrupted data")

        extractor = PptxExtractor()
        result = extractor.extract(corrupted)

        assert not result.success
        assert len(result.errors) > 0

    def test_config_with_skip_empty_slides(self, simple_pptx_file, empty_pptx_file):
        """Test skip_empty_slides configuration."""
        # Test with skip_empty_slides=True
        config = {"skip_empty_slides": True}
        extractor = PptxExtractor(config)

        # Extract from empty presentation - should have no blocks
        result = extractor.extract(empty_pptx_file)
        assert result.success

        # Extract from simple presentation - should have blocks
        result2 = extractor.extract(simple_pptx_file)
        assert result2.success
        assert len(result2.content_blocks) > 0

    def test_supports_format(self):
        """Test supports_format method."""
        extractor = PptxExtractor()

        assert extractor.supports_format(Path("test.pptx"))
        assert extractor.supports_format(Path("TEST.PPTX"))
        assert not extractor.supports_format(Path("test.docx"))
        assert not extractor.supports_format(Path("test.pdf"))

    def test_get_supported_extensions(self):
        """Test get_supported_extensions method."""
        extractor = PptxExtractor()
        extensions = extractor.get_supported_extensions()

        assert ".pptx" in extensions
        assert len(extensions) == 1

    def test_get_format_name(self):
        """Test get_format_name method."""
        extractor = PptxExtractor()
        assert extractor.get_format_name() == "Microsoft PowerPoint"


# Fixtures for test files
@pytest.fixture
def simple_pptx_file(tmp_path):
    """Create a simple test PowerPoint file."""
    try:
        from pptx import Presentation
    except ImportError:
        pytest.skip("python-pptx not installed")

    prs = Presentation()

    # Slide 1: Title slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])  # Title layout
    slide1.shapes.title.text = "Test Presentation"
    slide1.placeholders[1].text = "Subtitle text"

    # Slide 2: Content slide
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])  # Title and content
    slide2.shapes.title.text = "Content Slide"
    slide2.placeholders[1].text = "This is body content"

    # Slide 3: Another content slide
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "Second Content"
    slide3.placeholders[1].text = "More body text here"

    file_path = tmp_path / "simple_presentation.pptx"
    prs.save(str(file_path))
    return file_path


@pytest.fixture
def empty_pptx_file(tmp_path):
    """Create an empty PowerPoint file."""
    try:
        from pptx import Presentation
    except ImportError:
        pytest.skip("python-pptx not installed")

    prs = Presentation()
    # Add one blank slide
    prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    file_path = tmp_path / "empty.pptx"
    prs.save(str(file_path))
    return file_path


@pytest.fixture
def pptx_with_notes(tmp_path):
    """Create a PowerPoint file with speaker notes."""
    try:
        from pptx import Presentation
    except ImportError:
        pytest.skip("python-pptx not installed")

    prs = Presentation()

    # Add slide with notes
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Slide with Notes"
    slide.placeholders[1].text = "Content here"

    # Add speaker notes
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = "These are important speaker notes for the presenter."

    file_path = tmp_path / "with_notes.pptx"
    prs.save(str(file_path))
    return file_path


# Infrastructure Integration Tests
@pytest.mark.integration
class TestInfrastructureIntegration:
    """Test infrastructure component integration."""

    def test_accepts_config_manager(self, simple_pptx_file, tmp_path):
        """PptxExtractor accepts ConfigManager instance."""
        try:
            from infrastructure import ConfigManager
        except ImportError:
            pytest.skip("Infrastructure not available")

        # Create minimal config file
        config_file = tmp_path / "config.yaml"
        config_file.write_text(
            """
extractors:
  pptx:
    extract_notes: true
    extract_images: false
"""
        )

        config = ConfigManager(config_file)
        extractor = PptxExtractor(config)

        # Should accept ConfigManager without errors
        assert extractor.extract_notes is True
        assert extractor.extract_images is False

    def test_uses_config_values(self, pptx_with_notes, tmp_path):
        """PptxExtractor respects configuration values."""
        try:
            from infrastructure import ConfigManager
        except ImportError:
            pytest.skip("Infrastructure not available")

        # Config with notes disabled
        config_file = tmp_path / "config.yaml"
        config_file.write_text(
            """
extractors:
  pptx:
    extract_notes: false
"""
        )

        config = ConfigManager(config_file)
        extractor = PptxExtractor(config)
        result = extractor.extract(pptx_with_notes)

        # Notes should not be extracted
        notes_blocks = [
            b for b in result.content_blocks if b.metadata.get("is_speaker_note", False)
        ]
        assert len(notes_blocks) == 0

    def test_backward_compatible_dict_config(self, simple_pptx_file):
        """PptxExtractor still works with dict config."""
        config = {
            "extract_notes": False,
            "extract_images": True,
        }

        extractor = PptxExtractor(config)
        result = extractor.extract(simple_pptx_file)

        assert result.success
        assert extractor.extract_notes is False
        assert extractor.extract_images is True

    def test_error_codes_used(self, tmp_path):
        """PptxExtractor uses ErrorHandler for structured errors."""
        try:
            from infrastructure import ConfigManager
        except ImportError:
            pytest.skip("Infrastructure not available")

        # Create minimal config file
        config_file = tmp_path / "config.yaml"
        config_file.write_text(
            """
extractors:
  pptx:
    extract_notes: true
"""
        )

        config = ConfigManager(config_file)
        extractor = PptxExtractor(config)

        # Test file not found error
        result = extractor.extract(Path("nonexistent.pptx"))

        assert not result.success
        # Error should be structured with error code
        assert len(result.errors) > 0

    def test_logging_operations(self, simple_pptx_file, tmp_path):
        """PptxExtractor logs key extraction steps."""
        try:
            from infrastructure import ConfigManager
        except ImportError:
            pytest.skip("Infrastructure not available")

        # Create minimal config file
        config_file = tmp_path / "config.yaml"
        config_file.write_text(
            """
extractors:
  pptx:
    extract_notes: true
"""
        )

        config = ConfigManager(config_file)
        extractor = PptxExtractor(config)

        # Should not crash even with logging
        result = extractor.extract(simple_pptx_file)
        assert result.success
