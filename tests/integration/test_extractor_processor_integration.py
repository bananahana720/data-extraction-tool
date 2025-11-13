"""
Extractor → Processor Integration Tests.

Tests real data flow from extraction through processing stages:
- Extracted content → ContextLinker (hierarchy preservation)
- Extracted content → MetadataAggregator (metadata enrichment)
- Extracted content → QualityValidator (quality scoring)
- Cross-format processor compatibility
- Error propagation and graceful degradation
- Sequential processor chaining

Test IDs: EP-001 through EP-010
"""

import pytest

from extractors.excel_extractor import ExcelExtractor
from extractors.pptx_extractor import PptxExtractor
from src.core import ContentType, ProcessingStage
from src.extractors import DocxExtractor, PdfExtractor
from src.processors import ContextLinker, MetadataAggregator, QualityValidator

# ==============================================================================
# Test Markers
# ==============================================================================

pytestmark = [pytest.mark.integration, pytest.mark.extraction, pytest.mark.processing]


# ==============================================================================
# DOCX → Processor Integration Tests
# ==============================================================================


def test_ep_001_docx_to_context_linker_preserves_hierarchy(sample_docx_file):
    """
    Test EP-001: DOCX heading hierarchy preserved through ContextLinker.

    Scenario: DOCX with multi-level headings → Extract → ContextLinker

    Verifies:
    - Extraction successful with hierarchical content
    - ContextLinker processes without errors
    - Parent-child relationships established
    - Heading levels preserved in metadata
    """
    # Arrange: Extract DOCX with hierarchical content
    extractor = DocxExtractor()
    extraction_result = extractor.extract(sample_docx_file)

    # Assert: Extraction succeeded
    assert extraction_result.success is True
    assert len(extraction_result.content_blocks) > 0

    # Arrange: Process with ContextLinker
    processor = ContextLinker()

    # Act: Link context relationships
    processing_result = processor.process(extraction_result)

    # Assert: Processing succeeded
    assert processing_result.success is True
    assert len(processing_result.content_blocks) > 0

    # Assert: Hierarchy preserved
    # Find heading blocks
    headings = [
        block
        for block in processing_result.content_blocks
        if block.block_type == ContentType.HEADING
    ]

    if headings:
        # At least one heading should have hierarchy metadata
        # ContextLinker adds parent_id or hierarchy_level
        has_hierarchy = any(
            "parent_id" in block.metadata or "hierarchy_level" in block.metadata
            for block in headings
        )
        assert has_hierarchy, "ContextLinker should add hierarchy metadata"

    # Assert: All blocks have block_id (required for linking)
    for block in processing_result.content_blocks:
        assert block.block_id is not None


def test_ep_002_docx_to_metadata_aggregator_enriches_metadata(sample_docx_file):
    """
    Test EP-002: DOCX metadata enriched by MetadataAggregator.

    Scenario: DOCX → Extract → MetadataAggregator

    Verifies:
    - Document-level metadata preserved
    - Block-level metadata enriched
    - Word counts calculated
    - Statistics aggregated
    """
    # Arrange: Extract
    extractor = DocxExtractor()
    extraction_result = extractor.extract(sample_docx_file)

    assert extraction_result.success is True

    # Arrange: Process with MetadataAggregator
    processor = MetadataAggregator()

    # Act: Aggregate metadata
    processing_result = processor.process(extraction_result)

    # Assert: Processing succeeded
    assert processing_result.success is True

    # Assert: Document metadata preserved
    assert processing_result.document_metadata is not None
    assert processing_result.document_metadata.source_file == sample_docx_file

    # Assert: Processing stage set
    assert processing_result.processing_stage == ProcessingStage.METADATA_AGGREGATION

    # Assert: Block metadata enriched
    # MetadataAggregator adds word_count, char_count, etc.
    enriched_blocks = processing_result.content_blocks
    text_blocks = [b for b in enriched_blocks if b.block_type == ContentType.PARAGRAPH]

    if text_blocks:
        # At least some text blocks should have metadata
        has_enrichment = any(
            "word_count" in block.metadata or "length" in block.metadata for block in text_blocks
        )
        assert has_enrichment or len(text_blocks[0].content) == 0


def test_ep_003_docx_to_quality_validator_scores_extraction(sample_docx_file):
    """
    Test EP-003: DOCX extraction quality scored by QualityValidator.

    Scenario: DOCX → Extract → QualityValidator

    Verifies:
    - Quality score computed
    - Score in valid range (0-100)
    - Valid DOCX gets high score (>85)
    - Quality metadata added to blocks
    """
    # Arrange: Extract
    extractor = DocxExtractor()
    extraction_result = extractor.extract(sample_docx_file)

    assert extraction_result.success is True

    # Arrange: Process with QualityValidator
    processor = QualityValidator()

    # Act: Validate quality
    processing_result = processor.process(extraction_result)

    # Assert: Processing succeeded
    assert processing_result.success is True

    # Assert: Quality score present
    assert processing_result.quality_score is not None

    # Assert: Score in valid range
    score = processing_result.quality_score
    assert 0.0 <= score <= 100.0, f"Score {score} out of valid range"

    # Assert: Valid DOCX should score high
    assert score > 85.0, f"Expected high score for valid DOCX, got {score}"

    # Assert: Processing stage set
    assert processing_result.processing_stage == ProcessingStage.QUALITY_VALIDATION


# ==============================================================================
# PDF → Processor Integration Tests
# ==============================================================================


def test_ep_004_pdf_to_metadata_aggregator_preserves_page_numbers(sample_pdf_file):
    """
    Test EP-004: PDF page numbers preserved through MetadataAggregator.

    Scenario: PDF → Extract → MetadataAggregator

    Verifies:
    - Page numbers in extracted blocks
    - Page numbers preserved after processing
    - Multi-page document handled correctly
    """
    # Arrange: Extract PDF
    extractor = PdfExtractor()
    extraction_result = extractor.extract(sample_pdf_file)

    assert extraction_result.success is True
    assert len(extraction_result.content_blocks) > 0

    # Assert: Extracted blocks have page numbers
    for block in extraction_result.content_blocks:
        assert block.position is not None
        assert block.position.page >= 1

    # Arrange: Process with MetadataAggregator
    processor = MetadataAggregator()

    # Act: Process
    processing_result = processor.process(extraction_result)

    # Assert: Processing succeeded
    assert processing_result.success is True

    # Assert: Page numbers still present after processing
    for block in processing_result.content_blocks:
        assert block.position is not None
        assert block.position.page >= 1, "Page numbers should be preserved"


def test_ep_005_pdf_to_context_linker_handles_sequential_content(sample_pdf_file):
    """
    Test EP-005: PDF sequential content linked by ContextLinker.

    Scenario: PDF → Extract → ContextLinker

    Verifies:
    - Sequential blocks linked
    - Previous/next relationships established
    - Page boundaries respected
    """
    # Arrange: Extract PDF
    extractor = PdfExtractor()
    extraction_result = extractor.extract(sample_pdf_file)

    assert extraction_result.success is True

    # Arrange: Process with ContextLinker
    processor = ContextLinker()

    # Act: Link context
    processing_result = processor.process(extraction_result)

    # Assert: Processing succeeded
    assert processing_result.success is True
    assert len(processing_result.content_blocks) > 0

    # Assert: Blocks have context metadata
    # ContextLinker may add previous_id, next_id, or similar
    blocks = processing_result.content_blocks
    if len(blocks) > 1:
        # Check if linking metadata added
        has_context = any("previous_id" in b.metadata or "next_id" in b.metadata for b in blocks)
        # Or at minimum, all blocks should have IDs for potential linking
        all_have_ids = all(b.block_id is not None for b in blocks)
        assert has_context or all_have_ids


# ==============================================================================
# PPTX → Processor Integration Tests
# ==============================================================================


def test_ep_006_pptx_to_quality_validator_maintains_slide_context(tmp_path, sample_docx_file):
    """
    Test EP-006: PPTX slide context maintained through QualityValidator.

    Scenario: PPTX → Extract → QualityValidator

    Verifies:
    - Slide numbers preserved
    - Slide-level context maintained
    - Quality scored per slide
    """
    # Note: Using sample_docx_file as fallback since PPTX fixture
    # may not be in integration conftest. Testing with available fixture.

    # Create simple PPTX for testing
    try:
        from pptx import Presentation

        prs = Presentation()
        slide1 = prs.slides.add_slide(prs.slide_layouts[1])
        slide1.shapes.title.text = "Test Slide"
        slide1.shapes.placeholders[1].text = "Content for testing"

        pptx_file = tmp_path / "test.pptx"
        prs.save(str(pptx_file))

    except ImportError:
        pytest.skip("python-pptx not installed")

    # Arrange: Extract PPTX
    extractor = PptxExtractor()
    extraction_result = extractor.extract(pptx_file)

    assert extraction_result.success is True

    # Arrange: Process with QualityValidator
    processor = QualityValidator()

    # Act: Validate
    processing_result = processor.process(extraction_result)

    # Assert: Processing succeeded
    assert processing_result.success is True
    assert processing_result.quality_score is not None

    # Assert: Slide context preserved (page number = slide number)
    for block in processing_result.content_blocks:
        if block.position:
            assert block.position.page >= 1


# ==============================================================================
# XLSX → Processor Integration Tests
# ==============================================================================


def test_ep_007_xlsx_to_all_processors_handles_tabular_data(tmp_path):
    """
    Test EP-007: XLSX tabular data through all processors.

    Scenario: XLSX → Extract → ContextLinker → MetadataAggregator → QualityValidator

    Verifies:
    - Tabular data extracted correctly
    - All processors handle table content
    - Sheet context preserved throughout
    - Quality reflects data completeness
    """
    # Create simple XLSX for testing
    try:
        import openpyxl

        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = "Test Sheet"

        # Add data
        ws["A1"] = "Header1"
        ws["B1"] = "Header2"
        ws["A2"] = "Data1"
        ws["B2"] = "Data2"

        xlsx_file = tmp_path / "test.xlsx"
        wb.save(str(xlsx_file))

    except ImportError:
        pytest.skip("openpyxl not installed")

    # Arrange: Extract XLSX
    extractor = ExcelExtractor()
    extraction_result = extractor.extract(xlsx_file)

    assert extraction_result.success is True
    assert len(extraction_result.content_blocks) > 0

    # Act: Process through all processors
    blocks = extraction_result.content_blocks

    # ContextLinker
    linker = ContextLinker()
    linked_result = linker.process(extraction_result)
    assert linked_result.success is True

    # MetadataAggregator
    aggregator = MetadataAggregator()
    aggregated_result = aggregator.process(linked_result.content_blocks)
    assert aggregated_result.success is True

    # QualityValidator
    validator = QualityValidator()
    validated_result = validator.process(aggregated_result.content_blocks)
    assert validated_result.success is True

    # Assert: Final result has all enrichments
    assert validated_result.quality_score is not None
    assert validated_result.quality_score > 0.0


# ==============================================================================
# Cross-Format Processor Integration Tests
# ==============================================================================


def test_ep_008_multiple_extractors_same_processor(sample_docx_file, sample_pdf_file, tmp_path):
    """
    Test EP-008: Multiple format types through same processor.

    Scenario: DOCX + PDF → Extract both → Same ContextLinker instance

    Verifies:
    - Processor handles multiple formats
    - No cross-contamination between files
    - Processor state independent per call
    """
    # Arrange: Extract both formats
    docx_extractor = DocxExtractor()
    pdf_extractor = PdfExtractor()

    docx_result = docx_extractor.extract(sample_docx_file)
    pdf_result = pdf_extractor.extract(sample_pdf_file)

    assert docx_result.success is True
    assert pdf_result.success is True

    # Arrange: Single processor instance
    processor = ContextLinker()

    # Act: Process both
    docx_processing = processor.process(docx_result)
    pdf_processing = processor.process(pdf_result)

    # Assert: Both succeeded
    assert docx_processing.success is True
    assert pdf_processing.success is True

    # Assert: Results are independent
    assert len(docx_processing.content_blocks) > 0
    assert len(pdf_processing.content_blocks) > 0

    # Block IDs should be different (no cross-contamination)
    docx_ids = {b.block_id for b in docx_processing.content_blocks}
    pdf_ids = {b.block_id for b in pdf_processing.content_blocks}
    assert len(docx_ids.intersection(pdf_ids)) == 0, "Block IDs should not overlap"


# ==============================================================================
# Error Handling Integration Tests
# ==============================================================================


def test_ep_009_extraction_errors_handled_by_processors(corrupted_docx_file):
    """
    Test EP-009: Processors gracefully handle partial extraction.

    Scenario: Corrupted DOCX → Extraction (partial failure) → Processor

    Verifies:
    - Processors don't crash on empty/minimal input
    - Error information preserved
    - Graceful degradation
    """
    # Arrange: Extract corrupted file (may fail)
    extractor = DocxExtractor()
    extraction_result = extractor.extract(corrupted_docx_file)

    # May succeed with warnings or fail entirely
    # Arrange: Process whatever we got
    processor = MetadataAggregator()

    # Act: Process (even if extraction failed/partial)
    if extraction_result.content_blocks:
        processing_result = processor.process(extraction_result)

        # Assert: Processing completes (may have warnings)
        assert processing_result.success is True or len(processing_result.errors) > 0
    else:
        # Empty extraction - processor should handle empty input
        processing_result = processor.process(())
        assert processing_result.success is True  # Empty is valid


def test_ep_010_processor_chain_propagates_enrichments(sample_docx_file):
    """
    Test EP-010: Sequential processor chain propagates enrichments.

    Scenario: Extract → ContextLinker → MetadataAggregator → QualityValidator

    Verifies:
    - Each processor adds metadata
    - Previous enrichments preserved
    - Final result has all enrichments
    - No data loss through chain
    """
    # Arrange: Extract
    extractor = DocxExtractor()
    extraction_result = extractor.extract(sample_docx_file)

    assert extraction_result.success is True

    blocks = extraction_result.content_blocks
    initial_block_count = len(blocks)

    # Act: Process through chain
    # Stage 1: ContextLinker
    linker = ContextLinker()
    linked = linker.process(extraction_result)
    assert linked.success is True

    # Stage 2: MetadataAggregator (should preserve ContextLinker metadata)
    aggregator = MetadataAggregator()
    aggregated = aggregator.process(linked.content_blocks)
    assert aggregated.success is True

    # Stage 3: QualityValidator (should preserve all previous metadata)
    validator = QualityValidator()
    validated = validator.process(aggregated.content_blocks)
    assert validated.success is True

    # Assert: Block count preserved
    assert len(validated.content_blocks) == initial_block_count

    # Assert: Final result has quality score
    assert validated.quality_score is not None

    # Assert: Enrichments accumulated
    # Check a sample block for metadata accumulation
    if validated.content_blocks:
        sample_block = validated.content_blocks[0]
        assert sample_block.metadata is not None
        assert isinstance(sample_block.metadata, dict)

        # Should have some enrichment from processors
        # (Specific keys depend on processor implementations)
        assert len(sample_block.metadata) > 0


# ==============================================================================
# TXT Extractor Integration Tests (NEW)
# ==============================================================================


def test_ep_011_txt_to_processors_handles_simple_structure(sample_text_file):
    """
    Test EP-011: TXT simple structure through processors.

    Scenario: TXT → Extract → All Processors

    Verifies:
    - TXT extraction handled correctly
    - Processors work with simple text structure
    - Quality score reflects text simplicity
    """
    # Arrange: Extract TXT
    # TXT uses DocxExtractor per current implementation
    extractor = DocxExtractor()
    extraction_result = extractor.extract(sample_text_file)

    # TXT extraction should succeed
    assert extraction_result.success is True or len(extraction_result.warnings) > 0

    if not extraction_result.content_blocks:
        pytest.skip("TXT extraction returned no blocks")

    # Act: Process through all processors
    blocks = extraction_result.content_blocks

    linker = ContextLinker()
    linked = linker.process(extraction_result)
    assert linked.success is True

    aggregator = MetadataAggregator()
    aggregated = aggregator.process(linked.content_blocks)
    assert aggregated.success is True

    validator = QualityValidator()
    validated = validator.process(aggregated.content_blocks)
    assert validated.success is True

    # Assert: Quality score computed
    assert validated.quality_score is not None


# ==============================================================================
# Empty Input Handling
# ==============================================================================


def test_ep_012_processors_handle_empty_input():
    """
    Test EP-012: All processors handle empty input gracefully.

    Scenario: Empty block list → Each Processor

    Verifies:
    - No crash on empty input
    - Success=True for empty input (valid case)
    - Appropriate default values
    """
    empty_blocks = ()

    # ContextLinker
    linker = ContextLinker()
    result = linker.process(empty_blocks)
    assert result.success is True
    assert len(result.content_blocks) == 0

    # MetadataAggregator
    aggregator = MetadataAggregator()
    result = aggregator.process(empty_blocks)
    assert result.success is True
    assert len(result.content_blocks) == 0

    # QualityValidator
    validator = QualityValidator()
    result = validator.process(empty_blocks)
    assert result.success is True
    assert len(result.content_blocks) == 0
    # Quality score may be None or 0 for empty input
