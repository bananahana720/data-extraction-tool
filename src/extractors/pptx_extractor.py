"""
PowerPoint (PPTX) Extractor - Microsoft PowerPoint Presentation Extraction

Extracts content from PowerPoint presentations including:
- Slide text (titles and body)
- Speaker notes
- Slide sequence and layout
- Presentation metadata

Design:
- TDD implementation following BaseExtractor interface
- Uses python-pptx library for parsing
- Infrastructure integration (ConfigManager, logging, error handling)
"""

from pathlib import Path
from typing import Optional, Union
from datetime import datetime, timezone
import hashlib
import logging

try:
    from pptx import Presentation
    from pptx.exc import PackageNotFoundError
except ImportError:
    raise ImportError(
        "python-pptx is required for PptxExtractor. " "Install with: pip install python-pptx"
    )

import sys

sys.path.insert(0, str(Path(__file__).parent.parent))

from core import (
    BaseExtractor,
    ContentBlock,
    ContentType,
    DocumentMetadata,
    ExtractionResult,
    ImageMetadata,
    Position,
)

# Import infrastructure components
try:
    from infrastructure import (
        ConfigManager,
        get_logger,
        ErrorHandler,
    )

    INFRASTRUCTURE_AVAILABLE = True
except ImportError:
    INFRASTRUCTURE_AVAILABLE = False


class PptxExtractor(BaseExtractor):
    """
    Extracts content from Microsoft PowerPoint (.pptx) files.

    Extracts:
    - Slide titles and body text
    - Speaker notes
    - Slide sequence and positioning
    - Presentation metadata

    Example:
        >>> extractor = PptxExtractor()
        >>> result = extractor.extract(Path("presentation.pptx"))
        >>> if result.success:
        ...     for block in result.content_blocks:
        ...         print(f"Slide {block.position.slide}: {block.content}")
    """

    def __init__(self, config: Optional[Union[dict, object]] = None):
        """
        Initialize PPTX extractor with optional configuration.

        Args:
            config: Configuration options (dict or ConfigManager):
                - extract_notes: Extract speaker notes (default: True)
                - extract_images: Extract image metadata (default: True)
                - skip_empty_slides: Skip slides with no content (default: False)
        """
        super().__init__(config if isinstance(config, dict) or config is None else {})

        # Detect ConfigManager
        is_config_manager = (
            INFRASTRUCTURE_AVAILABLE
            and hasattr(config, "__class__")
            and config.__class__.__name__ == "ConfigManager"
        )
        self._config_manager = config if is_config_manager else None

        # Initialize infrastructure
        if INFRASTRUCTURE_AVAILABLE:
            self.logger = get_logger(__name__)
            self.error_handler = ErrorHandler()
        else:
            self.logger = logging.getLogger(__name__)
            self.error_handler = None

        # Load configuration
        if self._config_manager:
            extractor_config = self._config_manager.get_section("extractors.pptx", default={})
            extract_notes_val = extractor_config.get("extract_notes")
            self.extract_notes = extract_notes_val if extract_notes_val is not None else True
            extract_images_val = extractor_config.get("extract_images")
            self.extract_images = extract_images_val if extract_images_val is not None else True
            skip_empty_val = extractor_config.get("skip_empty_slides")
            self.skip_empty_slides = skip_empty_val if skip_empty_val is not None else False
        elif isinstance(config, dict):
            self.extract_notes = config.get("extract_notes", True)
            self.extract_images = config.get("extract_images", True)
            self.skip_empty_slides = config.get("skip_empty_slides", False)
        else:
            self.extract_notes = True
            self.extract_images = True
            self.skip_empty_slides = False

    def supports_format(self, file_path: Path) -> bool:
        """
        Check if file is a PPTX file.

        Args:
            file_path: Path to check

        Returns:
            True if file has .pptx extension
        """
        return file_path.suffix.lower() == ".pptx"

    def get_supported_extensions(self) -> list[str]:
        """Return supported file extensions."""
        return [".pptx"]

    def get_format_name(self) -> str:
        """Return human-readable format name."""
        return "Microsoft PowerPoint"

    def extract(self, file_path: Path) -> ExtractionResult:
        """
        Extract content from PPTX file.

        Strategy:
        1. Validate file exists and is accessible
        2. Open presentation with python-pptx
        3. Iterate through slides
        4. Extract text from shapes (title, body)
        5. Extract speaker notes if configured
        6. Generate presentation metadata
        7. Return structured result

        Args:
            file_path: Path to PPTX file

        Returns:
            ExtractionResult with content blocks and metadata

        Note:
            - Returns success=False for file-level errors
            - Returns partial results if some slides fail
            - Logs warnings for recoverable issues
        """
        import time

        start_time = time.time()

        # Log extraction start
        if INFRASTRUCTURE_AVAILABLE:
            self.logger.info("Starting PPTX extraction", extra={"file": str(file_path)})

        errors = []
        warnings = []
        content_blocks = []

        # Step 1: Validate file
        is_valid, validation_errors = self.validate_file(file_path)
        if not is_valid:
            if self.error_handler:
                error = self.error_handler.create_error("E001", file_path=str(file_path))
                errors.append(self.error_handler.format_for_user(error))
                if INFRASTRUCTURE_AVAILABLE:
                    self.logger.error(
                        "File validation failed",
                        extra={"file": str(file_path), "errors": validation_errors},
                    )
            else:
                errors.extend(validation_errors)

            return ExtractionResult(
                success=False,
                errors=tuple(errors),
                document_metadata=DocumentMetadata(
                    source_file=file_path,
                    file_format="pptx",
                ),
            )

        try:
            # Step 2: Open presentation
            try:
                prs = Presentation(file_path)
            except PackageNotFoundError as e:
                return ExtractionResult(
                    success=False,
                    errors=(f"Invalid PPTX file: {str(e)}",),
                    document_metadata=DocumentMetadata(
                        source_file=file_path,
                        file_format="pptx",
                    ),
                )
            except Exception as e:
                return ExtractionResult(
                    success=False,
                    errors=(f"Failed to open PPTX file: {str(e)}",),
                    document_metadata=DocumentMetadata(
                        source_file=file_path,
                        file_format="pptx",
                    ),
                )

            # Step 3: Extract slides
            sequence_index = 0
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_blocks = []

                # Extract text from shapes
                for shape in slide.shapes:
                    if not hasattr(shape, "text") or not shape.text:
                        continue

                    text = shape.text.strip()
                    if not text:
                        continue

                    # Detect if this is a title or body text
                    block_type = self._detect_shape_type(shape, slide)

                    # Create content block
                    block = ContentBlock(
                        block_type=block_type,
                        content=text,
                        raw_content=shape.text,
                        position=Position(slide=slide_num, sequence_index=sequence_index),
                        confidence=1.0,
                        metadata={
                            "slide_number": slide_num,
                            "shape_name": shape.name if hasattr(shape, "name") else None,
                        },
                    )
                    slide_blocks.append(block)
                    sequence_index += 1

                # Extract speaker notes if configured
                if self.extract_notes and slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        notes_block = ContentBlock(
                            block_type=ContentType.COMMENT,
                            content=notes_text,
                            position=Position(slide=slide_num, sequence_index=sequence_index),
                            confidence=1.0,
                            metadata={
                                "slide_number": slide_num,
                                "is_speaker_note": True,
                            },
                        )
                        slide_blocks.append(notes_block)
                        sequence_index += 1

                # Skip empty slides if configured
                if self.skip_empty_slides and not slide_blocks:
                    continue

                content_blocks.extend(slide_blocks)

            # Check if we got any content
            if not content_blocks:
                warnings.append("No content extracted from presentation")

            # Step 3.5: Extract images if configured
            images = []
            if self.extract_images:
                images = self._extract_image_metadata(prs)
                if INFRASTRUCTURE_AVAILABLE:
                    self.logger.debug(f"Extracted {len(images)} images")

            # Step 4: Generate presentation metadata
            doc_metadata = self._extract_presentation_metadata(file_path, prs)

            # Update statistics
            total_chars = sum(len(b.content) for b in content_blocks)
            total_words = sum(len(b.content.split()) for b in content_blocks)

            doc_metadata = DocumentMetadata(
                source_file=doc_metadata.source_file,
                file_format=doc_metadata.file_format,
                file_size_bytes=doc_metadata.file_size_bytes,
                file_hash=doc_metadata.file_hash,
                title=doc_metadata.title,
                author=doc_metadata.author,
                created_date=doc_metadata.created_date,
                modified_date=doc_metadata.modified_date,
                subject=doc_metadata.subject,
                keywords=doc_metadata.keywords,
                word_count=total_words,
                character_count=total_chars,
                page_count=len(prs.slides),  # Use slide count for presentations
                extracted_at=doc_metadata.extracted_at,
                extractor_version="0.1.0",
            )

            # Step 5: Log completion and return result
            duration = time.time() - start_time
            if INFRASTRUCTURE_AVAILABLE:
                self.logger.info(
                    "PPTX extraction complete",
                    extra={
                        "file": str(file_path),
                        "blocks": len(content_blocks),
                        "slides": len(prs.slides),
                        "duration_seconds": round(duration, 3),
                    },
                )

            return ExtractionResult(
                content_blocks=tuple(content_blocks),
                document_metadata=doc_metadata,
                images=tuple(images),
                success=True,
                warnings=tuple(warnings),
            )

        except PermissionError as e:
            if self.error_handler:
                error = self.error_handler.create_error(
                    "E500", file_path=str(file_path), original_exception=e
                )
                errors.append(self.error_handler.format_for_user(error))
            else:
                errors.append(f"Permission denied reading file: {file_path}")

            if INFRASTRUCTURE_AVAILABLE:
                self.logger.error("Permission denied", extra={"file": str(file_path)})

        except Exception as e:
            if self.error_handler:
                error = self.error_handler.create_error(
                    "E150", file_path=str(file_path), original_exception=e
                )
                errors.append(self.error_handler.format_for_user(error))
            else:
                errors.append(f"Unexpected error during extraction: {str(e)}")

            if INFRASTRUCTURE_AVAILABLE:
                self.logger.error(
                    "Extraction failed", extra={"file": str(file_path), "error": str(e)}
                )

        # Return failed result
        return ExtractionResult(
            success=False,
            errors=tuple(errors),
            warnings=tuple(warnings),
            document_metadata=DocumentMetadata(
                source_file=file_path,
                file_format="pptx",
            ),
        )

    def _detect_shape_type(self, shape, slide) -> ContentType:
        """
        Detect content type based on shape properties.

        Heuristic:
        - Shapes in title placeholders → HEADING
        - Other text shapes → PARAGRAPH

        Args:
            shape: python-pptx Shape object
            slide: python-pptx Slide object

        Returns:
            ContentType enum value
        """
        # Check if shape is title placeholder
        if hasattr(shape, "is_placeholder") and shape.is_placeholder:
            placeholder_format = shape.placeholder_format
            if placeholder_format.type == 1:  # Title placeholder type
                return ContentType.HEADING

        # Check if shape name suggests it's a title
        if hasattr(shape, "name") and "title" in shape.name.lower():
            return ContentType.HEADING

        # Default to paragraph
        return ContentType.PARAGRAPH

    def _extract_presentation_metadata(
        self, file_path: Path, prs: Presentation
    ) -> DocumentMetadata:
        """
        Extract presentation-level metadata from PPTX file.

        Extracts both file system metadata and embedded document properties.

        Args:
            file_path: Path to file
            prs: python-pptx Presentation object

        Returns:
            DocumentMetadata with available properties
        """
        # File system metadata
        file_stat = file_path.stat()
        file_size = file_stat.st_size

        # Generate file hash for deduplication
        file_hash = self._compute_file_hash(file_path)

        # Extract core properties (presentation metadata)
        core_props = prs.core_properties

        # Parse dates (may be None)
        created = core_props.created
        modified = core_props.modified

        # Parse keywords (may be None or string)
        keywords = ()
        if core_props.keywords:
            keywords = tuple(k.strip() for k in core_props.keywords.split(",") if k.strip())

        return DocumentMetadata(
            source_file=file_path,
            file_format="pptx",
            file_size_bytes=file_size,
            file_hash=file_hash,
            title=core_props.title or None,
            author=core_props.author or None,
            created_date=created if isinstance(created, datetime) else None,
            modified_date=modified if isinstance(modified, datetime) else None,
            subject=core_props.subject or None,
            keywords=keywords,
            extracted_at=datetime.now(timezone.utc),
        )

    def _compute_file_hash(self, file_path: Path) -> str:
        """
        Compute SHA256 hash of file for deduplication.

        Args:
            file_path: Path to file

        Returns:
            Hex string of SHA256 hash
        """
        sha256 = hashlib.sha256()

        # Read in chunks to handle large files
        with open(file_path, "rb") as f:
            for chunk in iter(lambda: f.read(8192), b""):
                sha256.update(chunk)

        return sha256.hexdigest()

    def _extract_image_metadata(self, prs: Presentation) -> list[ImageMetadata]:
        """
        Extract image metadata from presentation slides.

        Strategy:
        1. Iterate through all slides
        2. Find shapes that are pictures
        3. Extract image properties (dimensions, format)
        4. Create ImageMetadata objects

        Args:
            prs: python-pptx Presentation object

        Returns:
            List of ImageMetadata objects
        """
        from pptx.enum.shapes import MSO_SHAPE_TYPE

        images = []

        try:
            for slide_num, slide in enumerate(prs.slides, start=1):
                for shape in slide.shapes:
                    # Check if shape is a picture
                    if not hasattr(shape, "shape_type"):
                        continue

                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            # Get image object
                            image = shape.image

                            # Get dimensions (convert from EMUs to pixels)
                            # EMU (English Metric Unit): 914400 EMUs = 1 inch
                            # Assuming 96 DPI: 1 inch = 96 pixels
                            width_pixels = (
                                int(shape.width * 96 / 914400) if hasattr(shape, "width") else None
                            )
                            height_pixels = (
                                int(shape.height * 96 / 914400)
                                if hasattr(shape, "height")
                                else None
                            )

                            # Get image format
                            img_format = None
                            if hasattr(image, "ext"):
                                img_format = image.ext.upper()

                            # Create metadata
                            image_meta = ImageMetadata(
                                width=width_pixels,
                                height=height_pixels,
                                format=img_format,
                                # Store slide location in alt_text for now
                                alt_text=f"Image on slide {slide_num}",
                            )
                            images.append(image_meta)

                        except Exception as e:
                            # Log warning but continue with other images
                            if INFRASTRUCTURE_AVAILABLE:
                                self.logger.warning(
                                    f"Failed to extract image metadata from slide {slide_num}",
                                    extra={"error": str(e)},
                                )

        except Exception as e:
            # Log error but return partial results
            if INFRASTRUCTURE_AVAILABLE:
                self.logger.error("Error during image metadata extraction", extra={"error": str(e)})

        return images
